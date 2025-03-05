import json
import os
import re
import time

import pandas as pd
import pytesseract
import requests
import textract
from azure.ai.translation.text import TextTranslationClient
from azure.cognitiveservices.vision.computervision import ComputerVisionClient
from azure.core.credentials import AzureKeyCredential
from docx import Document
from msrest.authentication import CognitiveServicesCredentials
from nltk.corpus import stopwords
from nltk.stem import WordNetLemmatizer
from nltk.tokenize import word_tokenize
from openai import AzureOpenAI
from pdf2image import convert_from_path
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity


def azure_translate(text, source_language, dest_language):
    headers = {
        "Ocp-Apim-Subscription-Key": os.environ["AZURE_TRANSLATION_SUBSCRIPTION_KEY"],
        "Ocp-Apim-Subscription-Region": os.environ["AZURE_TRANSLATION_REGION"],
        "Content-type": "application/json",
    }
    url = "https://api.cognitive.microsofttranslator.com/translate?api-version=3.0&from={}&to={}".format(
        source_language, dest_language
    )
    body = [{"text": text}]
    response = requests.post(url, headers=headers, json=body)
    if response.status_code != 200:
        return "Error: the translation service failed." + str(response)
    return json.loads(response.content.decode("utf-8-sig"))


# text='Que la actividad actual de la Sociedad consiste en la gestión y alquiler de habitaciones de uso privativo en viviendas compartidas (en adelante, la "Actividad").Que la Sociedad actualmente está llevando adelante una ronda de financiación (la "Ronda de Financiación") a los fines de captar inversiones de terceros por hasta un importe máximo global de [250.000,00€] (en adelante, el "Importe Máximo de Financiación"). [Los fondos recaudados serán utilizados para continuar desarrollando la Actividad de la Sociedad]'
# translate(text, "es", "en")


def azure_ocr():
    # Set up your endpoint and key
    ocr_endpoint = os.environ["AZURE_COMPUTER_VISION_ENDPOINT"]
    ocr_api_key = os.environ["AZURE_COMPUTER_VISION_KEY"]

    # Create a client
    client = ComputerVisionClient(
        ocr_endpoint, CognitiveServicesCredentials(ocr_api_key)
    )

    # Path to the image or PDF file
    file_path = "/Users/terrylennon/Downloads/for_ocr.pdf"

    # Open the file and send it for analysis
    with open(file_path, "rb") as image_stream:
        ocr_result = client.recognize_printed_text_in_stream(
            image_stream, language="en"
        )

    # Get the OCR results
    for region in ocr_result.regions:
        for line in region.lines:
            line_text = " ".join([word.text for word in line.words])
            print(line_text)


def azure_translate_python():
    translation_endpoint = (
        "https://api.cognitive.microsofttranslator.com"  # Your endpoint URL
    )
    translation_api_key = "3QpZ7oQtWVOSLlUZTXD6p65kfpnL6hiKXVmr3SMnqL2E96aj4mfZJQQJ99AJACfhMk5XJ3w3AAAbACOGV6Ij"
    client = TextTranslationClient(
        endpoint=translation_endpoint,
        credential=AzureKeyCredential(translation_api_key),
    )

    client = TextTranslationClient(
        endpoint=translation_endpoint,
        credential=AzureKeyCredential(translation_api_key),
        region="swedencentral",
    )

    # Define the text to translate and the target language(s)
    text_to_translate = ["Hello, how are you?"]
    target_languages = ["es", "fr"]  # Translate to Spanish and French
    response = client.translate(body=text_to_translate, to_language=target_languages)

    for document in response:
        print(f"Original Text: {text_to_translate[0]}")
        for translation in document.translations:
            print(f"Translation ({translation.to}): {translation.text}")


""" OCR """


def ocr_with_textract(pdf_path):
    text = textract.process(pdf_path).decode("utf-8")
    return text


def ocr_with_pytesseract(pdf_path, page_limit=5):
    pytesseract.pytesseract.tesseract_cmd = r"/opt/homebrew/bin/tesseract"
    pages = convert_from_path(pdf_path)
    text = ""

    # Iterate through all the pages stored above. Restrict to 5 pages
    for page in pages[:page_limit]:
        # Use pytesseract to do OCR on the page
        text += pytesseract.image_to_string(page)
    return text


def ocr_with_ocrspace(pdf_path):
    return pdf_path


def read_docx(file_path):
    # Open the Word document
    doc = Document(file_path)

    # Initialize an empty string to store the text
    full_text = []

    # Content in tables
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text for cell in row.cells]
            full_text.append("\t".join(row_text))

    # Iterate through each paragraph in the document
    for para in doc.paragraphs:
        full_text.append(para.text)

    # Join the text in the list into a single string, separated by newlines
    return "\n".join(full_text)


def read_excel(file_path):
    # Read the Excel file
    df = pd.read_excel(file_path)

    # Initialize an empty string to store the text
    full_text = []

    # Iterate through each row in the DataFrame
    for index, row in df.iterrows():
        # Join the values of the row into a single string, separated by tabs
        row_text = "\t".join(map(str, row.values))
        full_text.append(row_text)

    # Join the text in the list into a single string, separated by newlines
    return "\n".join(full_text)


def read_pptx(file_path):
    # Load the PowerPoint file
    presentation = Presentation(file_path)

    # Initialize an empty string to store the text
    full_text = []

    # Iterate through each slide in the presentation
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)

    # Join the text in the list into a single string, separated by newlines
    return "\n".join(full_text)


def ocr_file(file_path, page_limit=5):
    if file_path.lower().endswith(".docx"):
        text = read_docx(file_path)
        return text

    if file_path.lower().endswith(".xlsx"):
        text = read_excel(file_path)
        return text

    if file_path.lower().endswith(".pptx"):
        text = read_pptx(file_path)
        return text

    if file_path.lower().endswith(".doc"):
        text = textract.process(file_path).decode("utf-8")
        return text

    if file_path.lower().endswith(".pdf"):
        text = textract.process(file_path).decode("utf-8")
        if len(text) < 1000:
            text = ocr_with_pytesseract(file_path, page_limit)
        return text


def process_docs(folder_path):
    outputs = []
    errors = []
    files = os.listdir(folder_path)
    for file in files:
        print(file)
        try:
            text = ocr_file(folder_path + "/" + file, 10)
            outputs.append({"file": file, "text": text})
        except:
            errors.append(file)
            pass
    return outputs


""" AI """


def process_xlsx_with_azure(
    xlsx_file_path, prompt, row_name, truncation=10000, themes=[""]
):
    df = pd.read_excel(xlsx_file_path)
    errors = []
    outputs = []
    for index, row in df.iterrows():
        print(row)
        try:
            analysis = analyse_with_azure(
                prompt, row[row_name], truncation, themes=[""]
            )
            row["analysis"] = analysis
            outputs.append(row)
        except:
            errors.append(row)
            pass
    return outputs, errors


def azure_openai_client():
    return AzureOpenAI(
        api_key=os.environ["AZURE_OPENAI_KEY"],
        api_version=os.environ["AZURE_OPENAI_VERSION"],
        azure_endpoint=os.environ["AZURE_OPENAI_ENDPOINT"],
    )


def analyse_with_azure(prompt, text, truncation=10000, themes=[""]):
    client = azure_openai_client()

    if len(themes) > 0:
        matches = find_all_paragraphs_with_themes(text, themes)
        matches_combined = " ".join(matches)
    else:
        matches_combined = text

    truncated = matches_combined[:truncation]

    messages = [
        {"role": "system", "content": prompt},
        {"role": "assistant", "content": truncated},
    ]

    response = client.chat.completions.create(
        messages=messages, model="gpt-4o-cdo", max_tokens=2000
    )
    time.sleep(1)

    print(response.choices[0].message.content)
    return response.choices[0].message.content.replace("```json", "").replace("```", "")


def translate_with_azure(language, text, truncation=100000):
    client = azure_openai_client()

    messages = [
        {"role": "system", "content": "Translate this to " + language + ": "},
        {"role": "assistant", "content": text[:truncation]},
    ]

    response = client.chat.completions.create(
        messages=messages, model="gpt-4o-cdo", max_tokens=2000
    )

    return response.choices[0].message.content


""" Utils """


def read_folder_file_paths(folder_path):
    file_paths = []
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        file_paths.append(file_path)

    return file_paths


def read_text_file(file_path):
    with open(file_path, "r") as file:
        contents = file.read()
        return contents


def find_all_paragraphs_with_themes(text, themes):
    matching_paragraphs = []  # Initialize an empty list to store matching paragraphs
    paragraphs = re.split("\n{2,}", text)
    for paragraph in paragraphs:
        if any(theme in paragraph for theme in themes):
            matching_paragraphs.append(
                paragraph
            )  # Add the paragraph to the list if it contains any of the themes
    return matching_paragraphs  # Return the list of matching paragraphs


def find_nearest_paragraphs_to_themes(text, themes):
    paragraphs = text.split("\n\n")
    vectorizer = TfidfVectorizer()
    tfidf_matrix = vectorizer.fit_transform(paragraphs + themes)
    similarity_scores = cosine_similarity(tfidf_matrix[-1], tfidf_matrix[:-1])
    threshold = 0.5
    relevant_paragraphs = [
        paragraphs[i]
        for i, score in enumerate(similarity_scores[0])
        if score > threshold
    ]

    return relevant_paragraphs


def create_dict(text_input):
    # Split the text into lines
    lines = text_input.strip().split("\n")

    # Initialize an empty dictionary
    data_dict = {}

    # Iterate over each line
    for line in lines:
        key, value = line.split(": ", 1)
        data_dict[key] = value

    return data_dict


def clean_content(text):
    no_line_breaks = text.replace("\n\n", " ").replace("\n", " ")
    pattern = r"[^a-zA-Z0-9 .]"  # Note the space and full stop characters after 0-9
    return re.sub(pattern, "", no_line_breaks)


def extract_if_numeric(x):
    substring = x[:8]  # Get the first 8 characters
    if substring.isdigit():  # Check if the substring is all digits
        return substring
    return None


def preprocess_text(text):
    # Initialize lemmatizer
    lemmatizer = WordNetLemmatizer()

    # Define stop words
    stop_words = set(stopwords.words("english"))

    # Remove URLs
    text = re.sub(r"http\S+", "", text)

    # Remove emails
    text = re.sub(r"\S+@\S+", "", text)

    # Remove punctuation and numbers
    text = re.sub(r"[^a-zA-Z\s]", "", text)

    # Tokenize and remove stop words
    tokens = word_tokenize(text)
    filtered_tokens = [word for word in tokens if word.lower() not in stop_words]

    # Lemmatize tokens
    lemmatized_tokens = [lemmatizer.lemmatize(word) for word in filtered_tokens]

    # Join tokens back to a string
    processed_text = " ".join(lemmatized_tokens)

    return processed_text


def azure_ocr():
    import os
    import time

    import requests
    from pdf2image import convert_from_path

    # Replace with your Azure Computer Vision subscription key and endpoint
    subscription_key = os.environ["AZURE_OCR_SUBSCRIPTION_KEY"]
    endpoint = os.environ["AZURE_OCR_ENDPOINT"]

    # API URL
    ocr_url = endpoint + "vision/v3.2/read/analyze"

    # Headers for the API request
    headers = {
        "Ocp-Apim-Subscription-Key": subscription_key,
        "Content-Type": "application/octet-stream",
    }

    # Path to your PDF file
    pdf_file = "original.pdf"

    # Convert PDF to images
    print("Converting PDF to images...")
    pages = convert_from_path(pdf_file, 300)  # 300 DPI for high resolution
    image_files = []

    for i, page in enumerate(pages):
        image_filename = f"page_{i + 1}.png"
        page.save(image_filename, "PNG")
        image_files.append(image_filename)
        print(f"Saved {image_filename}")

    # Perform OCR on each image
    print("\nPerforming OCR on images...")
    for image_path in image_files:
        print(f"\nProcessing {image_path}...")

        # Read the image into a byte array
        with open(image_path, "rb") as f:
            image_data = f.read()

        # Send the image data to the OCR API
        response = requests.post(ocr_url, headers=headers, data=image_data)
        response.raise_for_status()

        # Get the operation URL from the response headers
        operation_url = response.headers["Operation-Location"]

        # Poll the API until the OCR operation is complete
        analysis = {}
        poll = True
        while poll:
            response_final = requests.get(
                operation_url, headers={"Ocp-Apim-Subscription-Key": subscription_key}
            )
            analysis = response_final.json()
            status = analysis.get("status")
            print(f"Status: {status}")
            if status == "succeeded":
                poll = False
            elif status == "failed":
                print("OCR failed for this image.")
                poll = False
            else:
                time.sleep(1)

        # Extract and print the recognized text
        if status == "succeeded":
            print("OCR succeeded. Extracted text:")
            read_results = analysis["analyzeResult"]["readResults"]
            for page in read_results:
                for line in page["lines"]:
                    print(line["text"])
        else:
            print(f"No text extracted from {image_path}.")

    # Optional: Clean up the image files
    for image_path in image_files:
        os.remove(image_path)
        print(f"Deleted {image_path}")

    print("\nOCR process completed.")
